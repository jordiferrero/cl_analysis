import os, glob

from pptx import Presentation
from pptx.util import Inches

####
root = os.path.abspath(r'G:\My Drive\PhD\projects\external_measurements')
session = r'20220630_Alessandro_TUE_2D3D_CL'
folder = os.path.join(root, session)

hyp_maps_ending = "*/*processed.hspy"
spectras_file_ending = '*B.png'
####

# Create new PWP
prs = Presentation()

title_only_slide_layout = prs.slide_layouts[5]
blank_slide_layout = prs.slide_layouts[6]

# Iterate through HYP maps
hyp_files = glob.glob(os.path.join(folder, hyp_maps_ending))
hyp_folders = [os.path.dirname(f) for f in hyp_files]
hyp_folders.sort()

for hyp_folder in hyp_folders:
    slide = prs.slides.add_slide(blank_slide_layout)

    # Add title
    left = top = Inches(0)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.text = os.path.basename(hyp_folder)

    # Add bandpass image
    try:
        im = glob.glob(os.path.join(hyp_folder, "im_bandpass*"))[0]
        pic = slide.shapes.add_picture(im, left=Inches(0), top=Inches(1), width=Inches(10))
    except:
        pass

    # Add SE before
    try:
        im = glob.glob(os.path.join(hyp_folder, "SE_Before*"))[0]
        pic = slide.shapes.add_picture(im, left=Inches(5.8), top=Inches(4.5), width=Inches(2))
    except:
        pass

    # Add SE after
    try:
        im = glob.glob(os.path.join(hyp_folder, "SE_After*"))[0]
        pic = slide.shapes.add_picture(im, left=Inches(8), top=Inches(4.5), width=Inches(2))
    except:
        pass

    # Add spectra
    try:
        im = glob.glob(os.path.join(hyp_folder, "im_spectrum*"))[0]
        pic = slide.shapes.add_picture(im, left=Inches(0), top=Inches(4.5), width=Inches(5))
    except:
        pass

    # Add textbox
    txBox = slide.shapes.add_textbox(left=Inches(6), top=Inches(6.5), width=Inches(4), height=Inches(1))
    tf = txBox.text_frame
    tf.text = "Notes:"
    p = tf.add_paragraph()
    p.text = os.path.basename(hyp_folder)
    p.level = 1

# Iterate through SPECTRAS plots
spectras_path = os.path.join(folder, 'SPECTRAS', 'plots')
spectras_files = glob.glob(os.path.join(spectras_path, spectras_file_ending))
spectras_files.sort()

for spectras_f in spectras_files:
    slide = prs.slides.add_slide(title_only_slide_layout)

    # Add title
    title_shape = slide.shapes.title
    title_shape.text = os.path.basename(spectras_f).split('.')[0]

    # Add picture
    top = Inches(2)
    left = Inches(1)
    height = Inches(4.5)
    pic = slide.shapes.add_picture(spectras_f, left, top, height=height)

save_folder = os.path.join(folder, 'writing')
if not os.path.exists(save_folder):
    os.mkdir(save_folder)

prs.save(os.path.join(save_folder, f'results_{session}.pptx'))